"""PowerPoint generation using template files"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pathlib import Path
import logging
from typing import List, Dict, Optional

import config

logger = logging.getLogger(__name__)


class PPTGenerator:
    """Generate PowerPoint presentations from template files"""
    
    def __init__(self, template_name: str = None):
        """
        Initialize PPT Generator
        
        Args:
            template_name: Name of template file to use (e.g., 'professional', 'modern')
        """
        self.template_name = template_name or config.DEFAULT_TEMPLATE
        
        # Template file path
        self.template_path = Path(__file__).parent / 'templates' / f'{self.template_name}.pptx'
        
        if not self.template_path.exists():
            logger.warning(f"Template '{self.template_name}.pptx' not found, using default")
            self.template_path = Path(__file__).parent / 'templates' / f'{config.DEFAULT_TEMPLATE}.pptx'
        
        logger.info(f"PPT Generator initialized with template: {self.template_name}")
    
    def create_presentation(self, slides_data: List[Dict], output_path: str,
                           title: str = None, subtitle: str = None) -> str:
        """
        Create PowerPoint presentation from template
        
        Strategy: Modify existing template slides to preserve graphics/styling
        instead of deleting and recreating them.
        """
        logger.info(f"Creating presentation with {len(slides_data)} slides using {self.template_name} template")
        
        try:
            # Check if template file exists and load it
            if self.template_path.exists():
                prs = Presentation(str(self.template_path))
                logger.info(f"Loaded template from {self.template_path}")
                logger.info(f"📐 Template has {len(prs.slide_layouts)} layouts, {len(prs.slides)} example slides")
                
                # Log available layouts for debugging
                for i, layout in enumerate(prs.slide_layouts):
                    logger.debug(f"   Layout {i}: {layout.name}")
                
                # 1. Add title slide using Layout 0
                if len(prs.slide_layouts) > 0:
                    title_layout = prs.slide_layouts[0]
                    title_slide = prs.slides.add_slide(title_layout)
                    
                    # Fill title placeholders using idx, not array position
                    logger.debug(f"Creating title slide with title='{title}', subtitle='{subtitle}'")
                    
                    # Truncate subtitle if too long (max 120 chars for visual appeal)
                    if subtitle and len(subtitle) > 120:
                        subtitle = subtitle[:117] + "..."
                    
                    title_ph = self._get_placeholder_by_idx(title_slide, 0)
                    if title_ph:
                        title_ph.text = title if title else "Video Summary"
                        logger.info(f"✅ Filled title: '{title_ph.text}'")
                    else:
                        logger.error(f"❌ Could not find title placeholder (idx=0)")
                    
                    subtitle_ph = self._get_placeholder_by_idx(title_slide, 1)
                    if subtitle_ph:
                        subtitle_ph.text = subtitle if subtitle else "AI-Generated Summary"
                        logger.info(f"✅ Filled subtitle: '{subtitle_ph.text[:50]}...'")
                    else:
                        logger.error(f"❌ Could not find subtitle placeholder (idx=1)")
                    
                    logger.debug("Created title slide from Layout 0")
                else:
                    self._add_title_slide(prs, title, subtitle)
                
                
                # ═══════════════════════════════════════════════════════════
                # VALIDATION PIPELINE: Gemini → Validate → Map → Create
                # ═══════════════════════════════════════════════════════════
                
                logger.info("🔄 Starting validation pipeline...")
                
                # STAGE 1: Validate & Structure
                # Convert raw Gemini sections into validated slide data
                validated_slides = self._validate_and_structure_slides({'sections': slides_data})
                
                if not validated_slides:
                    logger.error("No valid slides after validation")
                    raise ValueError("Validation failed: No slides generated")
                
                # CREATE TABLE OF CONTENTS as Slide 2
                toc_slide = self._create_toc_slide(prs, validated_slides)
                logger.info(f"✅ Created Table of Contents slide (slide 2)")
                
                # STAGE 2: Map to Layouts
                # Determine which layout each slide should use
                logger.info(f"🗺️  Mapping {len(validated_slides)} slides to layouts...")
                for i, slide_data in enumerate(validated_slides):
                    layout_idx = self._map_slide_to_layout(slide_data, i)
                    slide_data['layout_idx'] = layout_idx
                
                # STAGE 3: Create Slides
                # Actually build the PowerPoint slides
                logger.info(f"🎨 Creating {len(validated_slides)} slides...")
                for i, slide_data in enumerate(validated_slides):
                    layout_idx = slide_data.get('layout_idx', 1)
                    
                    # Create slide from layout
                    if layout_idx < len(prs.slide_layouts):
                        layout = prs.slide_layouts[layout_idx]
                        new_slide = prs.slides.add_slide(layout)
                        
                        # Fill placeholders based on layout type
                        layout_type = slide_data.get('layout_type', 'regular')
                        
                        if layout_type == 'comparison' and layout_idx == 2:
                            # Use original title with tags for comparison detection
                            slide_data_with_tags = slide_data.copy()
                            slide_data_with_tags['title'] = slide_data.get('original_title', slide_data['title'])
                            self._fill_three_content_placeholders(new_slide, slide_data_with_tags)
                        elif layout_type == 'image' and layout_idx == 3:
                            self._fill_picture_layout(new_slide, slide_data)
                        else:
                            # Regular content (Layout 1 or 4)
                            self._fill_regular_placeholders(new_slide, slide_data)
                        
                        logger.debug(f"  ✓ Slide {i+1}/{len(validated_slides)}: '{slide_data['title']}'")
                    else:
                        # Fallback
                        logger.warning(f"Layout {layout_idx} not found, using fallback")
                        self._add_content_slide(prs, slide_data)
                
                logger.info(f"✅ Pipeline complete: {len(validated_slides)} slides created")

                # CREATE THANK YOU SLIDE as final slide using Layout 6
                thank_you_slide = self._create_thank_you_slide(prs)
                
                logger.info(f"✅ Final presentation has {len(prs.slides)} slides")
                
            else:
                # Fallback: create blank presentation
                logger.warning(f"Template file not found: {self.template_path}")
                logger.warning("Creating blank presentation instead")
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
                
                self._add_title_slide(prs, title, subtitle)
                logger.info(f"✅ Created title slide")
                
                # AUTOMATICALLY CREATE TABLE OF CONTENTS AS SLIDE 2
                # First, validate and structure slides even in fallback
                validated_slides = self._validate_and_structure_slides({'sections': slides_data})
                toc_slide = self._create_toc_slide(prs, validated_slides)
                logger.info(f"✅ Created Table of Contents slide (slide 2)")
                
                for slide_data in slides_data:
                    self._add_content_slide(prs, slide_data)
            
            # Save presentation
            Path(output_path).parent.mkdir(parents=True, exist_ok=True)
            prs.save(output_path)
            logger.info(f"✅ Presentation saved to {output_path}")
            
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise
    
    def _get_placeholder_by_idx(self, slide, idx: int):
        """
        Find placeholder by its idx property (not array position)
        
        CRITICAL: slide.placeholders[10] tries to access array index 10,
        but we need to find the placeholder WHERE placeholder_format.idx == 10
        
        Args:
            slide: The slide object
            idx: The placeholder idx to find (e.g., 0, 1, 2, 10, 11)
            
        Returns:
            Placeholder object if found, None otherwise
        """
        try:
            for placeholder in slide.placeholders:
                if hasattr(placeholder, 'placeholder_format'):
                    if placeholder.placeholder_format.idx == idx:
                        return placeholder
        except Exception as e:
            logger.debug(f"Error finding placeholder {idx}: {e}")
        return None
    
    def _validate_and_structure_slides(self, raw_sections: Dict) -> List[Dict]:
        """
        Validate and structure Gemini output into clean slide data
        
        Pipeline Stage 1: Parse & Validate
        - Extract slides from Gemini's text output
        - Validate required fields
        - Fix common AI mistakes
        - Normalize bullet counts
        
        Args:
            raw_sections: Dict with 'sections' from Gemini
            
        Returns:
            List of validated slide dicts with structure:
            {
                'title': str,
                'content': str (newline-separated bullets),
                'layout_type': 'regular' | 'comparison' | 'image',
                'keyframe_path': str or None,
                'validated': bool
            }
        """
        validated_slides = []
        
        sections = raw_sections.get('sections', [])
        if not sections:
            logger.warning("No sections found in Gemini output")
            return []
        
        logger.info(f"📝 Validating {len(sections)} slides from Gemini...")
        
        for i, section in enumerate(sections):
            try:
                # Extract basic fields
                title = section.get('title', f'Untitled Slide {i+1}').strip()
                content_list = section.get('content', [])
                
                # Combine content list into string
                if isinstance(content_list, list):
                    content = '\n'.join([str(c).strip() for c in content_list if c])
                else:
                    content = str(content_list).strip()
                
                # Detect layout type from title
                title_upper = title.upper()
                if '[COMPARISON]' in title_upper or '[FEATURES]' in title_upper:
                    layout_type = 'comparison'
                    # Fix bullet count for comparison slides (need 6 for 3 columns with 2 points each)
                    bullets = [b for b in content.split('\n') if b.strip()]
                    if len(bullets) != 6:
                        logger.info(f"  ⚠️  Slide {i+1}: Fixing bullet count {len(bullets)} → 6 for comparison (3 columns, 2 points each)")
                        if len(bullets) < 6:
                            bullets.extend([''] * (6 - len(bullets)))  # Pad
                        else:
                            bullets = bullets[:6]  # Truncate
                        content = '\n'.join(bullets)
                else:
                    layout_type = 'regular'
                    # Regular slides should have 4-6 bullets (max 6)
                    bullets = [b for b in content.split('\n') if b.strip()]
                    if len(bullets) > 6:
                        logger.info(f"  ⚠️  Slide {i+1}: Truncating {len(bullets)} bullets → 6 (max allowed)")
                        bullets = bullets[:6]
                        content = '\n'.join(bullets)
                    elif len(bullets) < 3:
                        logger.warning(f"  ⚠️  Slide {i+1}: Only {len(bullets)} bullets (min 3 recommended)")
                
                # Check for keyframe (will be added later in the pipeline)
                keyframe_path = section.get('keyframe_path', None)
                
                # CRITICAL: Prioritize comparison layout over image layout
                # If slide is marked as comparison/features, keep that layout even if it has a keyframe
                # Only use image layout if it's NOT a comparison slide
                if keyframe_path and layout_type != 'comparison':
                    layout_type = 'image'
                
                # Normalize title (remove tags for storage, we'll re-detect them later)
                clean_title = title.replace('[COMPARISON]', '').replace('[FEATURES]', '').strip()
                
                validated_slide = {
                    'title': clean_title,
                    'original_title': title,  # Keep original for tag detection
                    'content': content,
                    'layout_type': layout_type,
                    'keyframe_path': keyframe_path,
                    'validated': True
                }
                
                validated_slides.append(validated_slide)
                logger.debug(f"  ✓ Slide {i+1}: '{clean_title}' ({layout_type})")
                
            except Exception as e:
                logger.error(f"  ✗ Slide {i+1}: Validation error: {e}")
                # Add fallback slide
                validated_slides.append({
                    'title': f'Slide {i+1}',
                    'original_title': f'Slide {i+1}',
                    'content': 'Content unavailable',
                    'layout_type': 'regular',
                    'keyframe_path': None,
                    'validated': False
                })
        
        # 🚨 CRITICAL VALIDATION: Enforce 1-2 comparison slides (not more, not less)
        comparison_slides = [slide for slide in validated_slides if slide['layout_type'] == 'comparison']
        comparison_count = len(comparison_slides)
        
        if comparison_count == 0 and len(validated_slides) > 1:
            # No comparison slides - create ONE
            logger.warning("⚠️  NO COMPARISON SLIDE FOUND! Creating one automatically...")
            
            # Convert the second-to-last slide into a comparison slide
            # (avoid converting the last slide as it's usually a conclusion)
            target_idx = len(validated_slides) - 2 if len(validated_slides) > 2 else 0
            target_slide = validated_slides[target_idx]
            
            # Get existing bullets
            bullets = [b for b in target_slide['content'].split('\n') if b.strip()]
            
            # Ensure we have exactly 6 bullets for comparison layout
            if len(bullets) < 6:
                bullets.extend([f'Additional point {i+1}' for i in range(6 - len(bullets))])
            elif len(bullets) > 6:
                bullets = bullets[:6]
            
            # Update the slide to be a comparison slide
            target_slide['layout_type'] = 'comparison'
            target_slide['original_title'] = target_slide['title'] + ' [FEATURES]'
            target_slide['content'] = '\n'.join(bullets)
            
            logger.warning(f"✅ Converted slide '{target_slide['title']}' to FEATURES slide (Layout 2)")
            logger.info(f"✅ Final count: 1 comparison/features slide")
            
        elif comparison_count > 2:
            # Too many comparison slides - convert extras back to regular
            logger.warning(f"⚠️  TOO MANY COMPARISON SLIDES ({comparison_count})! Limiting to 2...")
            
            # Keep first 2, convert rest to regular
            comparison_indices = [i for i, slide in enumerate(validated_slides) if slide['layout_type'] == 'comparison']
            slides_to_convert = comparison_indices[2:]  # Convert slides beyond the first 2
            
            for idx in slides_to_convert:
                validated_slides[idx]['layout_type'] = 'regular'
                # Remove tags from original title
                validated_slides[idx]['original_title'] = validated_slides[idx]['title']
                logger.info(f"  ♻️  Converted slide '{validated_slides[idx]['title']}' back to REGULAR slide")
            
            logger.info(f"✅ Final count: 2 comparison/features slides (limited from {comparison_count})")
            
        else:
            # 1-2 comparison slides - perfect!
            logger.info(f"✅ Found {comparison_count} comparison/features slide(s) - PERFECT!")

        
        logger.info(f"✅ Validated {len(validated_slides)} slides successfully")
        return validated_slides
    
    def _map_slide_to_layout(self, slide_data: Dict, slide_index: int) -> int:
        """
        Map validated slide data to template layout index
        
        Pipeline Stage 2: Layout Mapping
        - Determine which layout (0-4) to use
        - Based on layout_type, keyframe presence, and index
        
        PRIORITY ORDER:
        1. Comparison/Features slides (Layout 2) - even if they have keyframes
        2. Image slides with keyframes (Layout 3)  
        3. Regular content slides (Layouts 1 & 4) - alternate for variety
        
        Args:
            slide_data: Validated slide dict
            slide_index: Index in slides array (for alternating)
            
        Returns:
            Layout index (0-4)
        """
        layout_type = slide_data.get('layout_type', 'regular')
        has_keyframe = bool(slide_data.get('keyframe_path'))
        
        # PRIORITY 1: Layout 2 - 3-column Comparisons/Features
        # These take precedence even if they have keyframes
        if layout_type == 'comparison':
            logger.debug(f"  Layout 2 (Comparison) ← Slide {slide_index+1}: comparison/features")
            return 2
        
        # PRIORITY 2: Layout 3 - Images/Keyframes  
        # Only use picture layout if NOT a comparison slide
        if has_keyframe or layout_type == 'image':
            logger.debug(f"  Layout 3 (Picture) ← Slide {slide_index+1}: keyframe")
            return 3
        
        # PRIORITY 3: Layouts 1 & 4 - Alternate for regular content
        layout_idx = 1 if slide_index % 2 == 0 else 4
        logger.debug(f"  Layout {layout_idx} (Regular) ← Slide {slide_index+1}: regular content")
        return layout_idx
    
    def _is_comparison_slide(self, slide_data: Dict) -> bool:
        """
        Detect if slide is a comparison/feature slide based on title tags
        """
        title = slide_data.get('title', '').upper()
        
        # Check for explicit tags
        if '[COMPARISON]' in title or '[FEATURES]' in title:
            return True
        
        # Check for comparison keywords in title
        comparison_keywords = [
            'VS', 'VERSUS', 'COMPARISON', 'COMPARE',
            'PROS AND CONS', 'ADVANTAGES', 'DISADVANTAGES',
            'BEFORE AND AFTER', 'FEATURES', 'BENEFITS'
        ]
        
        return any(keyword in title for keyword in comparison_keywords)
    
    def _map_to_template_slide(self, content_index: int, slide_data: Dict, total_slides: int) -> int:
        """
        Map content slide to appropriate layout (0-4)
        
        Template layouts:
        0: Title Slide
        1: Title and Content (regular content)
        2: Two Content (3-column comparisons/features)
        3: Picture with Caption (keyframes/images)
        4: Custom Layout (alternate regular content)
        
        Returns: Layout index (1-4 for content slides)
        """
        has_keyframe = bool(slide_data.get('keyframe_path'))
        is_comparison = self._is_comparison_slide(slide_data)
        
        # Layout 3: Always use for keyframe images
        if has_keyframe:
            return 3
        
        # Layout 2: Use for comparisons/features (3 columns)
        if is_comparison:
            return 2
        
        # Layouts 1 & 4: Alternate for regular content to reduce repetitiveness
        if content_index % 2 == 0:
            return 1  # Title and Content
        else:
            return 4  # Custom Layout
    
    def _add_title_slide(self, prs: Presentation, title: str = None, subtitle: str = None):
        """
        Add title slide using Layout 0
        
       Layout 0 placeholders:
        - [0]: Title
        - [1]: Subtitle
        """
        try:
            # Use Layout 0 for title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            
            # Fill placeholder 0 (Title)
            if len(slide.placeholders) > 0 and title:
                slide.placeholders[0].text = title or config.PPT_TITLE
            
            # Fill placeholder 1 (Subtitle)
            if len(slide.placeholders) > 1 and subtitle:
                slide.placeholders[1].text = subtitle or config.PPT_SUBTITLE
            
            logger.debug("Added title slide (Layout 0)")
            
        except Exception as e:
            logger.error(f"Error adding title slide: {e}")
            raise
    
    def _add_content_slide(self, prs: Presentation, slide_data: Dict):
        """
        Fallback method: Add content slide using Layout 3 (with image) or Layout 1 (without image)
        
        Layout 3 (with keyframe):
        - [0]: Title
        - [1]: Picture placeholder
        - [10]: Content/Text
        
        Layout 1 (without keyframe):
        - [0]: Title
        - [1]: Content/Text
        """
        has_keyframe = bool(slide_data.get('keyframe_path'))
        
        try:
            # Choose layout based on whether we have a keyframe
            if has_keyframe:
                # Use Layout 3 for images
                layout_idx = 3 if len(prs.slide_layouts) > 3 else 1
                slide_layout = prs.slide_layouts[layout_idx]
                logger.debug(f"Using Layout {layout_idx} (with image) for slide: {slide_data.get('title')}")
            else:
                # Use Layout 1 for regular content
                slide_layout = prs.slide_layouts[1]
                logger.debug(f"Using Layout 1 (text only) for slide: {slide_data.get('title')}")
            
            slide = prs.slides.add_slide(slide_layout)
            
            # Fill placeholder 0 (Title)
            if len(slide.placeholders) > 0:
                slide.placeholders[0].text = slide_data.get('title', 'Untitled Slide')
            
            # For Layout 3: Insert image in placeholder 1, then content in placeholder 10
            if has_keyframe and layout_idx == 3:
                # Insert keyframe image in placeholder 1
                keyframe_path = slide_data.get('keyframe_path')
                if keyframe_path and Path(keyframe_path).exists() and len(slide.placeholders) > 1:
                    try:
                        slide.placeholders[1].insert_picture(keyframe_path)
                        logger.debug(f"Inserted image: {keyframe_path}")
                    except Exception as e:
                        logger.warning(f"Could not insert image {keyframe_path}: {e}")
                
                # Content in placeholder 10
                content_placeholder_idx = 10
            else:
                # For Layout 1: Content in placeholder 1
                content_placeholder_idx = 1
            
            # Fill content placeholder with bullet points
            if len(slide.placeholders) > content_placeholder_idx:
                content = slide_data.get('content', '')
                self._add_bullet_points(slide.placeholders[content_placeholder_idx], content)
            
        except Exception as e:
            logger.error(f"Error adding content slide '{slide_data.get('title')}': {e}")
            raise
    
    def _add_bullet_points(self, placeholder, content: str):
        """
        Add bullet points to a placeholder text frame
        
        Args:
            placeholder: The placeholder shape
            content: Newline-separated bullet points
        """
        if not content:
            return
        
        try:
            text_frame = placeholder.text_frame
            text_frame.clear()  # Clear any existing content
            
            # Split content by newlines
            bullets = [line.strip() for line in content.split('\n') if line.strip()]
            
            # Add each bullet as a paragraph
            for i, bullet in enumerate(bullets):
                # Clean bullet markers if present
                bullet_text = bullet.lstrip('-•* ').strip()
                
                if i == 0:
                    # First paragraph already exists
                    p = text_frame.paragraphs[0]
                else:
                    # Add new paragraphs for subsequent bullets
                    p = text_frame.add_paragraph()
                
                p.text = bullet_text
                p.level = 0  # Top-level bullet
            
            logger.debug(f"Added {len(bullets)} bullet points")
            
        except Exception as e:
            logger.error(f"Error adding bullet points: {e}")
    
    def _update_title_slide(self, slide, title: str = None, subtitle: str = None):
        """
        Update existing title slide's content while preserving graphics
        
        Finds and updates text placeholders without removing background graphics
        """
        try:
            # Find and update text placeholders
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                # Check if it's actually a placeholder
                if not hasattr(shape, 'is_placeholder') or not shape.is_placeholder:
                    continue
                
                ph_type = shape.placeholder_format.type
                
                # Type 3 = CENTER_TITLE, Type 1 = TITLE
                if ph_type in [1, 3] and title:
                    shape.text = title or config.PPT_TITLE
                    logger.debug(f"Updated title placeholder (type {ph_type})")
                
                # Type 4 = SUBTITLE
                elif ph_type == 4 and subtitle:
                    shape.text = subtitle or config.PPT_SUBTITLE
                    logger.debug(f"Updated subtitle placeholder")
            
            logger.debug("Updated title slide")
            
        except Exception as e:
            logger.error(f"Error updating title slide: {e}")
            raise
    
    def _update_two_content_slide(self, slide, slide_data: Dict):
        """
        Legacy method: Update slide with Layout 2 for comparisons/features
        
        Splits bullets into three columns (2 bullets each)
        """
        try:
            # Remove [COMPARISON] or [FEATURES] tag from title
            title = slide_data.get('title', '')
            title = title.replace('[COMPARISON]', '').replace('[FEATURES]', '').strip()
            
            # Get bullets and split into three groups
            content = slide_data.get('content', '')
            bullets = [line.strip() for line in content.split('\n') if line.strip()]
            
            # Ensure we have 6 bullets for 3 columns (2 bullets each)
            while len(bullets) < 6:
                bullets.append('')  # Pad if needed
            if len(bullets) > 6:
                bullets = bullets[:6]  # Truncate to max 6
            
            left_bullets = bullets[:2]
            middle_bullets = bullets[2:4]
            right_bullets = bullets[4:6]
            
            left_content = '\n'.join(left_bullets)
            middle_content = '\n'.join(middle_bullets)
            right_content = '\n'.join(right_bullets)
            
            # Find placeholders
            text_placeholders = []
            for shape in slide.shapes:
                if not hasattr(shape, 'is_placeholder') or not shape.is_placeholder:
                    continue
                
                ph_type = shape.placeholder_format.type
                
                # Type 1 = TITLE
                if ph_type == 1:
                    shape.text = title
                    logger.debug(f"Updated three-content title")
                
                # Type 2 = BODY, Type 7 = TEXT - collect for content
                elif ph_type in [2, 7]:
                    text_placeholders.append(shape)
            
            # Fill the three text areas
            if len(text_placeholders) >= 1:
                self._add_bullet_points(text_placeholders[0], left_content)
                logger.debug(f"Filled left content ({len(left_bullets)} bullets)")
            
            if len(text_placeholders) >= 2:
                self._add_bullet_points(text_placeholders[1], middle_content)
                logger.debug(f"Filled middle content ({len(middle_bullets)} bullets)")
            
            if len(text_placeholders) >= 3:
                self._add_bullet_points(text_placeholders[2], right_content)
                logger.debug(f"Filled right content ({len(right_bullets)} bullets)")
            
            logger.debug(f"Updated three-content slide: {title}")
            
        except Exception as e:
            logger.error(f"Error updating two-content slide: {e}")
            raise
    def _create_toc_slide(self, prs: Presentation, validated_slides: List[Dict]) -> any:
        """
        Create Table of Contents slide as slide 2 using Layout 5
        
        Layout 5 (Professional): idx=0 (Title), idx=1 (Content)
        Layout 5 (Educational): idx=0 (Title), idx=10 (Content)
        
        Args:
            prs: Presentation object
            validated_slides: List of validated slide dicts
            
        Returns:
            TOC slide object
        """
        try:
            # Use Layout 5 for Table of Contents
            if len(prs.slide_layouts) > 5:
                toc_layout = prs.slide_layouts[5]
                toc_slide = prs.slides.add_slide(toc_layout)
                
                # Fill title in placeholder idx=0
                title_ph = self._get_placeholder_by_idx(toc_slide, 0)
                if title_ph:
                    title_ph.text = "Table of Contents"
                    logger.debug(f"Filled TOC title in placeholder idx=0")
                
                # Collect all slide titles (without numbers - template handles numbering)
                slide_titles = [slide['title'] for slide in validated_slides]
                
                # Try idx=1 first (professional), then idx=10 (educational)
                content_filled = False
                content_ph = self._get_placeholder_by_idx(toc_slide, 1)
                if content_ph:
                    content_ph.text_frame.clear()
                    # Add each title as a separate bullet point
                    for i, title in enumerate(slide_titles):
                        if i == 0:
                            p = content_ph.text_frame.paragraphs[0]
                        else:
                            p = content_ph.text_frame.add_paragraph()
                        p.text = title
                        p.level = 0
                    logger.debug(f"Filled TOC content in placeholder idx=1 (professional)")
                    content_filled = True
                else:
                    # Try educational template placeholder
                    content_ph = self._get_placeholder_by_idx(toc_slide, 10)
                    if content_ph:
                        content_ph.text_frame.clear()
                        # Add each title as a separate bullet point
                        for i, title in enumerate(slide_titles):
                            if i == 0:
                                p = content_ph.text_frame.paragraphs[0]
                            else:
                                p = content_ph.text_frame.add_paragraph()
                            p.text = title
                            p.level = 0
                        logger.debug(f"Filled TOC content in placeholder idx=10 (educational)")
                        content_filled = True
                
                if not content_filled:
                    logger.warning(f"Could not find TOC content placeholder (tried idx=1, 10)")
                
                logger.info(f"📋 Created Table of Contents with {len(slide_titles)} slides")
                return toc_slide
            else:
                logger.warning(f"Layout 5 not found in template, skipping TOC slide")
                return None
                
        except Exception as e:
            logger.error(f"Error creating TOC slide: {e}")
            return None
    
    def _create_thank_you_slide(self, prs: Presentation) -> any:
        """
        Create Thank You slide as final slide using Layout 6
        
        Layout 6: idx=0 (Title only)
        
        Args:
            prs: Presentation object
            
        Returns:
            Thank You slide object
        """
        try:
            # Use Layout 6 for Thank You slide  
            if len(prs.slide_layouts) > 6:
                thank_you_layout = prs.slide_layouts[6]
                thank_you_slide = prs.slides.add_slide(thank_you_layout)
                
                # Fill title with "Thank You"
                title_ph = self._get_placeholder_by_idx(thank_you_slide, 0)
                if title_ph:
                    title_ph.text = "Thank You"
                    logger.debug(f"Filled Thank You title in placeholder idx=0")
                else:
                    logger.warning(f"Could not find title placeholder in Layout 6")
                
                logger.info(f"🙏 Created Thank You slide")
                return thank_you_slide
            else:
                logger.warning(f"Layout 6 not found in template, skipping Thank You slide")
                return None
                
        except Exception as e:
            logger.error(f"Error creating Thank You slide: {e}")
            return None
    def _fill_regular_placeholders(self, slide, slide_data: Dict):
        """
        Fill placeholders in regular content slide (Layout 1 or Layout 4)
        
        Layout 1 placeholders: [0]=Title, [1]=Content
        Layout 4 placeholders: [0]=Title, [10]=Content
        """
        try:
            title = slide_data.get('title', 'Untitled')
            content = slide_data.get('content', '')
            
            # Fill title in placeholder with idx=0
            title_ph = self._get_placeholder_by_idx(slide, 0)
            if title_ph:
                title_ph.text = title
                logger.debug(f"Filled title in placeholder idx=0")
            
            # Try placeholder idx=1 first (Layout 1), then idx=10 (Layout 4)
            content_filled = False
            content_ph = self._get_placeholder_by_idx(slide, 1)
            if content_ph:
                # Layout 1: Content in placeholder idx=1
                content_ph.text_frame.clear()
                self._add_bullet_points(content_ph, content)
                logger.debug(f"Filled content in placeholder idx=1 (Layout 1)")
                content_filled = True
            else:
                # Try Layout 4 educational: Content in placeholder idx=10
                content_ph = self._get_placeholder_by_idx(slide, 10)
                if content_ph:
                    content_ph.text_frame.clear()
                    self._add_bullet_points(content_ph, content)
                    logger.debug(f"Filled content in placeholder idx=10 (Layout 4)")
                    content_filled = True
                else:
                    # Try Layout 4 professional: Content in placeholder idx=2
                    content_ph = self._get_placeholder_by_idx(slide, 2)
                    if content_ph:
                        content_ph.text_frame.clear()
                        self._add_bullet_points(content_ph, content)
                        logger.debug(f"Filled content in placeholder idx=2 (Layout 4 professional)")
                        content_filled = True
            
            if not content_filled:
                logger.warning(f"Could not find content placeholder for regular slide")
            
            logger.debug(f"Filled regular slide: {title}")
            
        except Exception as e:
            logger.error(f"Error filling regular placeholders: {e}")
            raise
    
    
    def _fill_picture_layout(self, slide, slide_data: Dict):
        """
        Fill placeholders in Layout 3 (Picture with Caption)
        
        Layout 3 placeholders: [1]=Picture, [0]=Title, [10 or 2 or 11]=Content/Caption
        """
        try:
            title = slide_data.get('title', 'Untitled')
            content = slide_data.get('content', '')
            keyframe_path = slide_data.get('keyframe_path')
            
            # Log all available placeholders for debugging
            logger.info(f"📸 Picture layout for '{title}' - Available placeholders:")
            for ph in slide.placeholders:
                if hasattr(ph, 'placeholder_format'):
                    logger.info(f"  idx={ph.placeholder_format.idx}, type={ph.placeholder_format.type}, name={ph.name if hasattr(ph, 'name') else 'N/A'}")
            
            # Fill title in placeholder idx=0
            title_ph = self._get_placeholder_by_idx(slide, 0)
            if title_ph:
                title_ph.text = title
                logger.info(f"✅ Filled title in placeholder idx=0: '{title}'")
            else:
                logger.warning(f"❌ Could not find title placeholder (idx=0)")
            
            # Insert picture in placeholder idx=1
            if keyframe_path and Path(keyframe_path).exists():
                picture_ph = self._get_placeholder_by_idx(slide, 1)
                if picture_ph:
                    try:
                        picture_ph.insert_picture(keyframe_path)
                        logger.info(f"✅ Inserted image in placeholder idx=1: {Path(keyframe_path).name}")
                    except Exception as e:
                        logger.warning(f"❌ Could not insert image: {e}")
                else:
                    logger.warning(f"❌ Could not find picture placeholder (idx=1)")
            else:
                logger.info(f"ℹ️  No keyframe to insert (path={keyframe_path})")
            
            # Fill content/caption - try multiple placeholder indices
            # Different templates use different indices: 10 (educational), 2 (professional), 11 (some templates), 1 (fallback)
            content_filled = False
            
            for ph_idx in [10, 2, 11]:  # Try in order of likelihood
                content_ph = self._get_placeholder_by_idx(slide, ph_idx)
                if content_ph and hasattr(content_ph, 'text_frame'):
                    try:
                        content_ph.text_frame.clear()
                        self._add_bullet_points(content_ph, content)
                        logger.info(f"✅ Filled caption in placeholder idx={ph_idx} ({len(content.split(chr(10)))} bullets)")
                        content_filled = True
                        break
                    except Exception as e:
                        logger.debug(f"Could not use placeholder idx={ph_idx}: {e}")
            
            if not content_filled:
                logger.warning(f"⚠️  Could not find caption placeholder (tried idx=10, 2, 11)")
                logger.warning(f"⚠️  Content NOT inserted: {content[:50]}...")
            
            logger.debug(f"Filled picture layout slide: {title}")
            
        except Exception as e:
            logger.error(f"Error filling picture layout: {e}")
            raise

    
    def _fill_three_content_placeholders(self, slide, slide_data: Dict):
        """
        Fill placeholders in Layout 2 (Two Content with 3 columns)
        
        Layout 2 placeholders: [0]=Title, [2]=Left, [10]=Middle, [11]=Right
        Splits bullets into 3 columns (2 points each) for comparisons/features
        """
        try:
            # Remove tags from title
            title = slide_data.get('title', '')
            title = title.replace('[COMPARISON]', '').replace('[FEATURES]', '').strip()
            
            # Split content into three groups
            content = slide_data.get('content', '')
            bullets = [line.strip() for line in content.split('\n') if line.strip()]
            
            # Pad or truncate to ensure we have 6 bullets for 3 columns (2 points each)
            while len(bullets) < 6:
                bullets.append('')  # Pad if needed
            if len(bullets) > 6:
                bullets = bullets[:6]  # Truncate to max 6 bullets (2 per column)
            
            # Split into 3 groups (2 bullets each)
            left_bullets = bullets[:2]
            middle_bullets = bullets[2:4]
            right_bullets = bullets[4:6]
            
            left_content = '\n'.join(left_bullets)
            middle_content = '\n'.join(middle_bullets)
            right_content = '\n'.join(right_bullets)
            
            # Fill title in placeholder idx=0
            title_ph = self._get_placeholder_by_idx(slide, 0)
            if title_ph:
                title_ph.text = title
                logger.debug(f"Filled title in placeholder idx=0")
            
            # Fill left content in placeholder idx=2
            left_ph = self._get_placeholder_by_idx(slide, 2)
            if left_ph:
                left_ph.text_frame.clear()
                self._add_bullet_points(left_ph, left_content)
                logger.debug(f"Filled left column in placeholder idx=2")
            
            # Fill middle content in placeholder idx=10
            middle_ph = self._get_placeholder_by_idx(slide, 10)
            if middle_ph:
                middle_ph.text_frame.clear()
                self._add_bullet_points(middle_ph, middle_content)
                logger.debug(f"Filled middle column in placeholder idx=10")
            
            # Fill right content in placeholder idx=11
            right_ph = self._get_placeholder_by_idx(slide, 11)
            if right_ph:
                right_ph.text_frame.clear()
                self._add_bullet_points(right_ph, right_content)
                logger.debug(f"Filled right column in placeholder idx=11")
            
            logger.debug(f"Filled three-content slide: {title}")
            
        except Exception as e:
            logger.error(f"Error filling three-content placeholders: {e}")
            raise
    
    def _update_content_slide(self, slide, slide_data: Dict, prs: Presentation):
        """
        Update existing content slide while preserving graphics
        
        Updates text and optionally inserts image
        """
        has_keyframe = bool(slide_data.get('keyframe_path'))
        
        try:
            # Find and update placeholders
            for shape in slide.shapes:
                # Skip non-placeholders
                if not hasattr(shape, 'is_placeholder') or not shape.is_placeholder:
                    continue
                
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
                
                # Type 1 = TITLE
                if ph_type == 1:
                    shape.text = slide_data.get('title', 'Untitled Slide')
                    logger.debug(f"Updated title")
                
                # Type 18 = PICTURE
                elif ph_type == 18 and has_keyframe:
                    keyframe_path = slide_data.get('keyframe_path')
                    if keyframe_path and Path(keyframe_path).exists():
                        try:
                            shape.insert_picture(keyframe_path)
                            logger.debug(f"Inserted image: {keyframe_path}")
                        except Exception as e:
                            logger.warning(f"Could not insert image: {e}")
                
                # Type 2 = BODY, Type 7 = TEXT/OBJECT
                elif ph_type in [2, 7]:
                    content = slide_data.get('content', '')
                    self._add_bullet_points(shape, content)
            
            logger.debug(f"Updated content slide: {slide_data.get('title')}")
            
        except Exception as e:
            logger.error(f"Error updating content slide: {e}")
            raise


def get_available_templates():
    """Get list of available template files"""
    templates_dir = Path(__file__).parent / 'templates'
    templates_dir.mkdir(exist_ok=True)
    
    template_files = list(templates_dir.glob('*.pptx'))
    
    templates = []
    for template_file in template_files:
        template_name = template_file.stem
        templates.append({
            'id': template_name,
            'name': template_name.replace('_', ' ').title(),
            'description': f'{template_name} template',
            'preview': f'{template_name}.png'
        })
    
    # If no templates found, return default structure
    if not templates:
        templates.append({
            'id': 'professional',
            'name': 'Professional',
            'description': 'Clean professional template',
            'preview': 'professional.png'
        })
    
    return templates